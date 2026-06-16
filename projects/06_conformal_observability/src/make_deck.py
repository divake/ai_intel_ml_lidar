"""Build the CogniSense presentation deck (PPTX). Story: Problem -> Demo -> Solution -> Results.
No result shown twice (each as EITHER a plot OR a table). Demo slide embeds the two videos.
Output: results/CogniSense_deck.pptx
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys; sys.path.insert(0, os.path.dirname(__file__))
import common as C

FIG = C.FIGS; RES = C.RESULTS
INK = RGBColor(0x21, 0x25, 0x2b); MUT = RGBColor(0x5a, 0x63, 0x70)
ACC = RGBColor(0x1f, 0x5f, 0xb4); RED = RGBColor(0xc0, 0x39, 0x2b); WHITE = RGBColor(0xff, 0xff, 0xff)
FONT = "Times New Roman"
EMUIN = 914400

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of paragraphs; each paragraph = list of (txt,size,bold,color,italic)."""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (txt, size, bold, color, *ital) in para:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.italic = bool(ital and ital[0])
    return tb


def bullets(s, l, t, w, h, items, size=18, color=INK, gap=8):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # support (text) or (text, bold-lead)
        lead, rest = (it if isinstance(it, tuple) else (None, it))
        r0 = p.add_run(); r0.text = "•  "; r0.font.name = FONT; r0.font.size = Pt(size); r0.font.color.rgb = ACC; r0.font.bold = True
        if lead:
            rl = p.add_run(); rl.text = lead + "  "; rl.font.name = FONT; rl.font.size = Pt(size); rl.font.bold = True; rl.font.color.rgb = color
        rr = p.add_run(); rr.text = rest; rr.font.name = FONT; rr.font.size = Pt(size); rr.font.color.rgb = color
    return tb


def title(s, t, kicker=None):
    box(s, 0, 0, SW, Inches(1.15), fill=RGBColor(0xf4, 0xf6, 0xf9))
    box(s, 0, Inches(1.15), SW, Pt(3), fill=ACC)
    text(s, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.9),
         [[(t, 30, True, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, Inches(0.57), Inches(0.74), Inches(12.2), Inches(0.4),
             [[(kicker, 14, False, MUT, True)]])


def pic_fit(s, path, l, t, w, h):
    iw, ih = Image.open(path).size
    scale = min(w/iw, h/ih); pw, ph = int(iw*scale), int(ih*scale)
    return s.shapes.add_picture(path, int(l+(w-pw)/2), int(t+(h-ph)/2), width=pw, height=ph)


def footer(s, n):
    text(s, Inches(0.55), Inches(7.06), Inches(10), Inches(0.35),
         [[("Conformal Safety Bounds for LiDAR Localization", 10, False, MUT)]])
    text(s, Inches(12.3), Inches(7.06), Inches(0.7), Inches(0.35),
         [[(str(n), 10, False, MUT)]], align=PP_ALIGN.RIGHT)


# ============================================================ 1 · TITLE
s = slide()
box(s, 0, 0, SW, SH, fill=RGBColor(0x12, 0x1a, 0x2b))
box(s, 0, Inches(3.95), SW, Pt(3), fill=ACC)
text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.6),
     [[("Conformal Safety Bounds for LiDAR Localization", 40, True, WHITE)]])
text(s, Inches(0.92), Inches(4.15), Inches(11.5), Inches(0.8),
     [[("Teaching a self-driving robot to know when it can't trust where it is", 20, False, RGBColor(0xc8, 0xd2, 0xe0), True)]])
text(s, Inches(0.92), Inches(6.2), Inches(11.5), Inches(0.8),
     [[("Divake Kumar   ·   AEON Lab, UIC   ·   DARPA JUMP 2.0 CogniSense", 15, False, RGBColor(0xa9, 0xb6, 0xc8))]])

# ============================================================ 2 · PROBLEM
s = slide(); title(s, "The Problem — a robot that doesn't know it's lost", "Localization without honest uncertainty is unsafe")
bullets(s, Inches(0.55), Inches(1.55), Inches(6.6), Inches(5.2), [
    ("Point estimate.", "A LiDAR robot localizes by matching scans to a map and outputs a single best-guess pose."),
    ("No honest 'how sure'.", "In self-similar corridors the geometry is ambiguous — the estimate can drift or jump silently."),
    ("The filter lies.", "Its on-board confidence (EKF covariance) looks tiny and trustworthy — but it isn't (next results)."),
    ("What we need.", "A localization uncertainty the robot can act on: with a guarantee, in real time, on the edge."),
], size=19, gap=14)
pic_fit(s, os.path.join(FIG, "F5_map.png"), Inches(7.4), Inches(1.7), Inches(5.5), Inches(4.4))
text(s, Inches(7.4), Inches(6.15), Inches(5.5), Inches(0.6),
     [[("The 250 m corridor loop the robot drives autonomously.", 13, False, MUT, True)]], align=PP_ALIGN.CENTER)
footer(s, 2)

# ============================================================ 3 · DEMO  (videos)
s = slide(); title(s, "Demo — a self-driving LiDAR robot", "Real robot · real sensor · real building")
vids = [("results/robot_moving_40s.mp4", "_poster_robot.png", "Robot driving the corridor autonomously (phone capture)"),
        ("results/lidar_wide_timelapse.mp4", "_poster_lidar.png", "Live LiDAR mapping — FAST-LIO2 LiDAR-inertial SLAM")]
xs = [Inches(0.55), Inches(6.95)]
for (vid, poster, cap), x in zip(vids, xs):
    pth = os.path.join(C.PROJ, vid); pp = os.path.join(RES, poster)
    bw, bh = Inches(5.85), Inches(3.9)
    try:
        iw, ih = Image.open(pp).size; sc = min(bw/iw, bh/ih); pw, ph = int(iw*sc), int(ih*sc)
        cx, cy = int(x + (bw-pw)/2), int(Inches(1.7) + (bh-ph)/2)
        s.shapes.add_movie(pth, cx, cy, pw, ph, poster_frame_image=pp, mime_type="video/mp4")
    except Exception as e:
        box(s, x, Inches(1.7), bw, bh, fill=RGBColor(0xe9, 0xed, 0xf2), line=MUT)
        text(s, x, Inches(3.3), bw, Inches(0.6), [[("▶  insert video: " + os.path.basename(vid), 14, True, MUT)]], align=PP_ALIGN.CENTER)
    text(s, x, Inches(5.75), bw, Inches(0.7), [[(cap, 14, False, INK, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.55), Inches(6.5), Inches(12.25), Inches(0.5), fill=RGBColor(0xf4, 0xf6, 0xf9))
text(s, Inches(0.7), Inches(6.52), Inches(12.0), Inches(0.45),
     [[("RoboSense Helios-16 LiDAR + RealSense IMU  ·  FAST-LIO2 SLAM  ·  251.6 m loop  ·  801k-point loop-closed map  ·  runs on an Intel NUC", 14, True, ACC)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ============================================================ 4 · SOLUTION
s = slide(); title(s, "Solution — a calibrated safety bound from what the robot sees", "Distribution-free · online · no Bayesian priors")
steps = [("1 · Observe", "sensing quality the robot already has (scan density / effective range)"),
         ("2 · Bound", "predict the localization error with conformal prediction — a distribution-free guarantee"),
         ("3 · Adapt", "keep coverage valid online with Adaptive Conformal Inference (ACI), no retraining"),
         ("4 · Act", "govern speed — slow / stop when the bound grows; full speed when it's small")]
x = Inches(0.55); w = Inches(2.95); gap = Inches(0.18)
for i, (h, d) in enumerate(steps):
    lx = Emu(int(x) + i*(int(w)+int(gap)))
    box(s, lx, Inches(2.0), w, Inches(2.7), fill=RGBColor(0xf4, 0xf6, 0xf9), line=RGBColor(0xd5, 0xdd, 0xe6))
    box(s, lx, Inches(2.0), w, Inches(0.62), fill=ACC)
    text(s, lx, Inches(2.02), w, Inches(0.6), [[(h, 18, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(lx)+90000), Inches(2.8), Emu(int(w)-180000), Inches(1.8), [[(d, 15, False, INK)]])
    if i < 3:
        text(s, Emu(int(lx)+int(w)-30000), Inches(3.0), gap, Inches(0.6), [[("→", 22, True, ACC)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.55), Inches(5.2), Inches(12.3), Inches(1.2),
     [[("The contribution: turn ", 19, False, INK), ("“how sure am I”", 19, True, RED),
       (" into a number with a coverage guarantee — and into a ", 19, False, INK), ("safe action", 19, True, RED), (".", 19, False, INK)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# ============================================================ 5 · RESULT 1  (plot F2)
s = slide(); title(s, "Result 1 — the filter's confidence is not a safety bound", "Why you can't trust the EKF covariance")
pic_fit(s, os.path.join(FIG, "F2_overconfidence.png"), Inches(0.5), Inches(1.5), Inches(7.6), Inches(5.3))
bullets(s, Inches(8.25), Inches(2.0), Inches(4.7), Inches(4.5), [
    ("1.3%.", "Its nominal “90%” interval actually contains the true pose only ~1.3% of the time."),
    ("~9× over-confident.", "It claims millimeter precision while real errors reach tens of centimetres."),
    ("So:", "safety cannot be gated on the localizer's own covariance."),
], size=18, gap=16)
footer(s, 5)

# ============================================================ 6 · RESULT 2  (plot F1)
s = slide(); title(s, "Result 2 — reliability collapses as sensing degrades", "Less sensing → larger, directional error")
pic_fit(s, os.path.join(FIG, "F1_degradation_curve.png"), Inches(0.5), Inches(1.5), Inches(7.6), Inches(5.3))
bullets(s, Inches(8.25), Inches(2.0), Inches(4.7), Inches(4.5), [
    ("4× error.", "As effective range drops 40 → 3 m (sparse / occluded / cheaper sensors), error grows ~4×."),
    ("Failures triple.", "8% → 23% of scans become unreliable."),
    ("The aperture.", "Error along the corridor grows faster than across it — the geometry itself goes blind."),
], size=18, gap=16)
footer(s, 6)

# ============================================================ 7 · RESULT 3  (plot F3)
s = slide(); title(s, "Result 3 — conformal prediction gives an honest, online bound", "A guarantee that holds while the robot drives")
pic_fit(s, os.path.join(FIG, "F3_calibration.png"), Inches(0.5), Inches(1.55), Inches(8.4), Inches(5.1))
bullets(s, Inches(9.05), Inches(2.0), Inches(3.9), Inches(4.5), [
    ("Honest.", "A naive bound silently under-covers under shift (0.88 vs 0.90 target)."),
    ("Self-correcting.", "Adaptive Conformal Inference restores coverage online (→ 0.90), no retraining."),
    ("Distribution-free.", "no Bayesian priors, no model of the noise."),
], size=17, gap=14)
footer(s, 7)

# ============================================================ 8 · RESULT 4  (TABLE)
s = slide(); title(s, "Result 4 — it slows down exactly when it should", "Risk-aware speed: 70% fewer unsafe events")
rowsd = [("LiDAR range", "gate slows", "unsafe — gated", "unsafe — ungated"),
         ("40 m (good)", "10%", "316", "317"),
         ("8 m", "46%", "344", "462"),
         ("6 m", "77%", "119", "583"),
         ("5 m", "92%", "6", "573"),
         ("4 m", "100%", "0", "696"),
         ("3 m (poor)", "100%", "0", "863"),
         ("Total", "—", "1148", "3886")]
nr, nc = len(rowsd), 4
gt = s.shapes.add_table(nr, nc, Inches(0.7), Inches(1.7), Inches(7.6), Inches(4.7)).table
gt.columns[0].width = Inches(2.3)
for j in range(1, 4): gt.columns[j].width = Inches(1.77)
for i, row in enumerate(rowsd):
    for j, val in enumerate(row):
        cell = gt.cell(i, j); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val; r.font.name = FONT; r.font.size = Pt(15)
        head = (i == 0); tot = (i == nr-1)
        r.font.bold = head or tot
        r.font.color.rgb = WHITE if head else (RED if (tot and j == 2) else INK)
        if head: cell.fill.solid(); cell.fill.fore_color.rgb = ACC
        elif tot: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xf0, 0xe6, 0xe4)
        else: cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xf4, 0xf6, 0xf9)
bullets(s, Inches(8.65), Inches(2.1), Inches(4.3), Inches(4.5), [
    ("Selective.", "full speed when sensing is good; stop only when it's blind."),
    ("70% fewer unsafe events.", "3886 → 1148; zero in the worst conditions."),
    ("The filter can't do this.", "it never knows it's degraded."),
], size=17, gap=16)
footer(s, 8)

# ============================================================ 9 · TAKEAWAYS
s = slide(); title(s, "Takeaways", "What to remember")
bullets(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(4.6), [
    ("A real autonomous LiDAR robot —", "and an uncertainty it can act on, not just a pose."),
    ("The on-board filter is overconfident", "(its “90%” is right 1.3% of the time); conformal prediction gives the calibrated, online bound it cannot."),
    ("Reliability degrades predictably with sensing —", "and the robot slows exactly when it should: 70% fewer unsafe events."),
    ("Distribution-free, no Bayesian priors, edge-ready", "— a safety layer for any LiDAR localizer."),
    ("Next:", "close the loop on the live robot; richer degradation & multi-environment data."),
], size=20, gap=18)
box(s, Inches(0.7), Inches(6.35), Inches(11.9), Pt(2.5), fill=ACC)
footer(s, 9)

out = os.path.join(RES, "CogniSense_deck.pptx")
prs.save(out)
print("wrote", out, "(", round(os.path.getsize(out)/1e6, 1), "MB,", len(prs.slides._sldIdLst), "slides )")
